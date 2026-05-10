from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Function to add a slide
def add_slide(prs, layout_idx, title_text, bullet_points):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    if len(bullet_points) > 0 and hasattr(slide.shapes, 'placeholders') and len(slide.shapes.placeholders) > 1:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = bullet_points[0]
        
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
    return slide

# Slide 1: Title
title_slide_layout = 0
slide = prs.slides.add_slide(prs.slide_layouts[title_slide_layout])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Project Elevate: The Intelligence Terminal for Healthcare"
subtitle.text = "Strategic Partnership & Co-Founder Roadmap\n\nIsaac Wang, Technical Founder\n[Candidate Name], Prospective Business Co-Founder"

# Slide 2: Vision
add_slide(prs, 1, "The $2 Billion Commercialization Gap", [
    "The Bottleneck: Brilliant healthcare innovations fail because researchers lack clear market and regulatory data.",
    "The Fragmentation: There are over 46,733 federal health signals across the FDA, CDC, CMS, and Census—but they are siloed and inaccessible.",
    "Our Solution: We’ve built an intelligence engine that synthesizes this data into a 45-second, fully validated strategic brief."
])

# Slide 3: Moat
add_slide(prs, 1, "An Unfair Advantage From Day One", [
    "Live Infrastructure: The MVP is fully built using a robust full-stack and PostgreSQL architecture.",
    "The Data Lake: We currently have 46,733 normalized federal demand signals embedded and indexed for semantic search.",
    "Multi-Agent Validation: We use a LangGraph 5-node pipeline where an independent AI 'Critic' audits math and regulatory claims to ensure absolute accuracy."
])

# Slide 4: Market
add_slide(prs, 1, "Market Potential & Exit Landscape", [
    "Total Addressable Market: $4.2B global market for healthcare R&D intelligence and regulatory advisory.",
    "Target Customers: 6,200+ federally funded researchers, 300+ University Tech Transfer Offices (TTOs), and mid-sized biotech BD teams.",
    "Market Comparables: B2B healthcare intelligence platforms (like Citeline or IQVIA subsidiaries) trade at 10x-15x revenue multiples.",
    "The Goal: Build the definitive 'Predictive Innovation Intent' dataset and position for a strategic acquisition within 3–5 years."
])

# Slide 5: Funding Strategy
add_slide(prs, 1, "How We Fund Our Growth", [
    "The Dilemma: We need capital to run the AI, but VCs want to see paying users. Users want a functioning product before paying.",
    "Step 1: Non-Dilutive Capital (Now). Leverage university ecosystems and tech grants to cover our API/server costs for 6–12 months.",
    "Step 2: Scrappy Traction (Next). Secure our first 15-20 paying users to prove product-market fit.",
    "Step 3: Venture Capital (Soon). Use that initial revenue to pitch Y Combinator and early-stage funds from a position of power."
])

# Slide 6: Step 1 (Runway)
add_slide(prs, 1, "Immediate Goal: Secure the Runway (Weeks 1-3)", [
    "AWS Activate / Google Cloud for Startups: Apply using university credentials to secure up to $100k in compute credits.",
    "OpenAI/Anthropic Startup Programs: Apply for API credits to cover embedding and report-generation costs.",
    "Skandalaris Venture Competition: Draft the executive summary and pitch our live MVP for up to $25k in non-dilutive capital.",
    "WashU LEAP Inventor Challenge: Position our platform as a tool to accelerate academic IP to market for a $50k grant opportunity."
])

# Slide 7: Step 2 (Customers)
add_slide(prs, 1, "Getting Our First 15 Customers (Months 1-3)", [
    "Target MRR: $1,500",
    "1. Sourcing: Monitor NIH RePORTER for newly awarded grants in our core domains (Oncology, Cardiology, etc.).",
    "2. The 'Give-First' Email: Offer a custom-generated Market Access & FDA pathway brief specifically for their newly funded project.",
    "3. The Demo: Run their innovation through the platform live on a 15-minute Zoom call to show the AI catching regulatory edge cases.",
    "4. The Frictionless Close: Convert them at $89/month. This falls under discretionary spending for lab P-cards."
])

# Slide 8: Step 3 (Scaling)
add_slide(prs, 1, "Moving Upmarket & Y Combinator (Months 4-12)", [
    "Target MRR: $15k - $30k+",
    "Y Combinator: With 15+ paying users, apply to YC's Bio/Health track to secure $500k and build engineering capacity.",
    "University Tech Transfer Offices (TTOs): Pitch institutional site licenses ($1,500/month) to save them hours vetting invention disclosures.",
    "Mid-Sized Biotech: Target early-stage commercial biotech companies looking to optimize their trial sites and regulatory pathways."
])

# Slide 9: Step 4 (Exit)
add_slide(prs, 1, "Becoming the Industry Standard (Year 2+)", [
    "Target MRR: $100k - $250k+",
    "Raise Series A ($3M - $5M): Scale our enterprise sales team and expand the database internationally (EMA, PMDA).",
    "Pharma BD Subscriptions: License our aggregate data to Big Pharma. Show them what top researchers are building before it hits the market.",
    "Continuous Engagement: Roll out live surveillance RSS alerts so users stay subscribed for competitive intelligence."
])

# Slide 10: Partnership
add_slide(prs, 1, "A Balanced Partnership", [
    "Isaac’s Focus (Product & Engineering):",
    "- Maintain and scale the PostgreSQL/vector architecture.",
    "- Expand the dataset past 100,000 signals.",
    "- Ensure 99.9% uptime and preserve the 45-second latency.",
    "Your Focus (Revenue & Growth):",
    "- Execute the micro-grant strategy in Week 1.",
    "- Run the 4-step acquisition playbook to secure our first paying cohort.",
    "- Navigate institutional sales and investor pitching (Y Combinator)."
])

# Save the presentation
prs.save('/Users/isaacwang/Downloads/ProjectElevate/Project_Elevate_Cofounder_Pitch.pptx')
print("✅ PowerPoint file created successfully: ~/Downloads/ProjectElevate/Project_Elevate_Cofounder_Pitch.pptx")

