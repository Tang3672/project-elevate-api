from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Function to add a slide
def add_slide(prs, layout_idx, title_text, bullet_points):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    if len(bullet_points) > 0 and hasattr(slide.shapes, 'placeholders') and len(slide.shapes.placeholders) > 1:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = bullet_points[0]
        
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
    return slide

# Slide 1: Title
title_slide_layout = 0
slide = prs.slides.add_slide(prs.slide_layouts[title_slide_layout])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "PROJECT ELEVATE"
subtitle.text = "Partnership Proposal & Execution Mandate\n\nIsaac Wang, Technical Founder\n[Candidate Name], Business Co-Founder"

# Slide 2: Equity & Exit
add_slide(prs, 1, "Equity, Exit, & Personal Wealth", [
    "This is a B2B intelligence terminal designed for acquisition within 36-48 months.",
    "1. Near-Term (Months 6-12): Market Salary. Upon closing our $500k Pre-Seed, we institute market-rate founder salaries.",
    "2. Mid-Term (Months 18-24): Secondary Liquidity. During Series A, we negotiate secondary shares to de-risk personally.",
    "3. Strategic Exit (Year 3-4): $10M+ Payout. B2B healthcare SaaS commands high multiples. Exit targeted at $50M-$100M."
])

# Slide 3: Solving the Chicken & Egg
add_slide(prs, 1, "The Funding Sequence", [
    "The Dilemma: We need capital to run the AI, but VCs want paying users. Users want a functioning product before paying.",
    "1. Now - Stop the Bleed (Micro-Grants): Acquire $5k - $50k in compute credits and university grants for free API usage.",
    "2. Next - Scrappy Traction (First 15 Customers): We sell the data, not the UI. 15 users at $89/mo = $1,335 MRR.",
    "3. Then - Pre-Seed / Y Combinator ($500k): With a live product and paying users, we pitch YC and angels from a position of power."
])

# Slide 4: Week 1 Priority
add_slide(prs, 1, "Week 1 Priority: Non-Dilutive Capital", [
    "AWS Activate / Google Cloud: Apply using our university domain to zero out hosting and database costs.",
    "Skandalaris Venture Competition: Draft executive summary. Leverage our live MVP as proof of technical execution.",
    "WashU LEAP Inventor Challenge: Submit translational research application to accelerate academic IP.",
    "Arch Grants (St. Louis): Pitch as St. Louis-built B2B infrastructure for $75,000 equity-free cash."
])

# Slide 5: Customer Acquisition
add_slide(prs, 1, "The Acquisition Playbook", [
    "We are running highly-targeted, data-driven sniper outreach to individuals with fresh funding.",
    "1. Sourcing: Filter NIH RePORTER for 'Phase I SBIR' or 'R01' grants awarded in the last 30 days.",
    "2. The 'Give-First' Email: Send a customized FDA pathway & Market Access brief tailored to their project.",
    "3. The 'Aha!' Demo: Generate a report live on a 15-minute Zoom to show the AI catching regulatory edge cases.",
    "4. The Frictionless Close: Convert at $89/month via Stripe so they can charge it directly to a lab P-Card."
])

# Slide 6: The Moat
add_slide(prs, 1, "Why They Won't 'Just Build It Themselves'", [
    "1. The Data Infrastructure: We maintain 10 asynchronous API connectors embedding 46,733 signals into a PostgreSQL pgvector database.",
    "2. Multi-Agent Validation: We use a LangGraph Pipeline where an independent 'Critic AI' audits math and regulatory claims before the user sees it.",
    "3. Our edge: We sell certainty and validated research, not just text generation."
])

# Slide 7: Timeline & MRR Targets
add_slide(prs, 1, "Growth Trajectory & MRR Targets", [
    "Months 1-2 (Survival & Seeding): Secure Skandalaris & AWS Credits. Acquire first 15 paid users. Target: $1.5k MRR.",
    "Months 3-5 (Pre-Seed & YC): Apply to Y Combinator. Raise $500k to hire engineers. Target: $500k Raised.",
    "Months 6-12 (Mid-Market & TTOs): Pitch Tech Transfer Offices ($1,500/mo) and mid-sized biotechs. Target: $30k+ MRR.",
    "Months 18+ (Series A & Exit): Raise $3M-$5M Series A. Position for strategic acquisition. Target: $250k+ MRR."
])

# Slide 8: The Mandate
add_slide(prs, 1, "A Balanced Partnership", [
    "Isaac’s Domain (Product & Tech):",
    "- Scale the MoE LangGraph architecture and expand the database (100k+ signals).",
    "- Build Enterprise API and ensure 45-second latency.",
    "Your Domain (Revenue & Growth):",
    "- Week 1: Secure compute credits and submit Skandalaris grants.",
    "- Month 1: Acquire the first 15 paying customers via the playbook.",
    "- Month 3: Prepare Y Combinator / Pre-Seed pitch deck.",
    "- Month 6: Close first B2B Tech Transfer Office contract."
])

# Save the presentation
prs.save('/Users/isaacwang/Downloads/ProjectElevate/Project_Elevate_Updated_Pitch.pptx')
print("✅ PowerPoint file created successfully: ~/Downloads/ProjectElevate/Project_Elevate_Updated_Pitch.pptx")

